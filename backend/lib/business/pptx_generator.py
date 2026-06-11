"""TOOL 6 — PowerPoint Generator. One reusable MG&CO branded template:
dark `#131313` title slide with white/blue typography, clean white content
slides with a blue `#2d7ff9` accent bar, and a 'Powered by Jarvis OS1' footer.
"""
import io

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

from backend.lib.business.connectors.base import ConnectorResult
from backend.lib.business.document_store import save_document
from backend.lib.business.model_router import SONNET
from backend.lib.business.real_estate.llm import call_claude, parse_json_response

DARK = RGBColor(0x13, 0x13, 0x13)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
BLUE = RGBColor(0x2D, 0x7F, 0xF9)
GRAY = RGBColor(0x88, 0x88, 0x88)
BODY_TEXT = RGBColor(0x33, 0x33, 0x33)

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)
FONT = "Arial"

SLIDE_PLANS: dict[str, list[str]] = {
    "listing": ["Property Hero & Stats", "Key Features", "Neighborhood", "Pricing", "Next Steps"],
    "cma": ["Comparable Sales", "Pricing Strategy", "Timeline"],
    "buyer_guide": ["The Buying Process", "What to Expect", "Financing Basics", "Next Steps"],
    "custom": [],
}

DECK_LABELS: dict[str, str] = {
    "listing": "Listing Presentation",
    "cma": "Comparative Market Analysis",
    "buyer_guide": "Buyer's Guide",
    "custom": "Presentation",
}

_PLAN_SYSTEM = "You write concise, persuasive real estate presentation content for a slide deck. Output JSON only."


def _set_text(text_frame, text, size, color, bold=False, align=PP_ALIGN.LEFT):
    text_frame.word_wrap = True
    p = text_frame.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.color.rgb = color
    run.font.bold = bold
    run.font.name = FONT


def _add_footer(slide, dark_bg=False):
    box = slide.shapes.add_textbox(Inches(0.5), SLIDE_H - Inches(0.4), Inches(6), Inches(0.3))
    color = RGBColor(0x66, 0x66, 0x66) if dark_bg else GRAY
    _set_text(box.text_frame, "Powered by Jarvis OS1", 9, color)


def _add_accent_bar(slide):
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(0.18), SLIDE_H)
    bar.fill.solid()
    bar.fill.fore_color.rgb = BLUE
    bar.line.fill.background()


def _blank_slide(prs):
    return prs.slides.add_slide(prs.slide_layouts[6])


def _add_title_slide(prs, title, subtitle):
    slide = _blank_slide(prs)

    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_W, SLIDE_H)
    bg.fill.solid()
    bg.fill.fore_color.rgb = DARK
    bg.line.fill.background()

    accent = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, Inches(3.6), Inches(0.18), Inches(1.6))
    accent.fill.solid()
    accent.fill.fore_color.rgb = BLUE
    accent.line.fill.background()

    brand_box = slide.shapes.add_textbox(Inches(0.9), Inches(0.6), Inches(6), Inches(0.4))
    _set_text(brand_box.text_frame, "JARVIS OS1 · MG&CO", 12, BLUE, bold=True)

    title_box = slide.shapes.add_textbox(Inches(0.9), Inches(2.9), Inches(11.5), Inches(1.5))
    _set_text(title_box.text_frame, title, 40, WHITE, bold=True)

    if subtitle:
        sub_box = slide.shapes.add_textbox(Inches(0.9), Inches(4.1), Inches(11.5), Inches(0.8))
        _set_text(sub_box.text_frame, subtitle, 18, RGBColor(0x9A, 0xA0, 0xA6))

    _add_footer(slide, dark_bg=True)
    return slide


def _add_content_slide(prs, heading, bullets):
    slide = _blank_slide(prs)
    _add_accent_bar(slide)

    head_box = slide.shapes.add_textbox(Inches(0.7), Inches(0.5), Inches(12), Inches(0.9))
    _set_text(head_box.text_frame, heading or "", 28, DARK, bold=True)

    body_box = slide.shapes.add_textbox(Inches(0.7), Inches(1.6), Inches(12), Inches(5.2))
    tf = body_box.text_frame
    tf.word_wrap = True
    for i, bullet in enumerate(bullets or ["—"]):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        run = p.add_run()
        run.text = f"•  {bullet}"
        run.font.size = Pt(16)
        run.font.color.rgb = BODY_TEXT
        run.font.name = FONT
        p.space_after = Pt(10)

    _add_footer(slide)
    return slide


def _add_table_slide(prs, heading, headers, rows):
    slide = _blank_slide(prs)
    _add_accent_bar(slide)

    head_box = slide.shapes.add_textbox(Inches(0.7), Inches(0.5), Inches(12), Inches(0.9))
    _set_text(head_box.text_frame, heading or "", 28, DARK, bold=True)

    headers = headers or ["Item"]
    rows = rows or []
    n_rows = len(rows) + 1
    n_cols = len(headers)
    table_shape = slide.shapes.add_table(n_rows, n_cols, Inches(0.7), Inches(1.6), Inches(12), Inches(0.6 * n_rows))
    table = table_shape.table

    for c, h in enumerate(headers):
        cell = table.cell(0, c)
        cell.text = str(h)
        cell.fill.solid()
        cell.fill.fore_color.rgb = DARK
        for p in cell.text_frame.paragraphs:
            for run in p.runs:
                run.font.color.rgb = WHITE
                run.font.bold = True
                run.font.size = Pt(13)
                run.font.name = FONT

    for r, row in enumerate(rows, start=1):
        for c in range(n_cols):
            val = row[c] if c < len(row) else ""
            cell = table.cell(r, c)
            cell.text = str(val)
            for p in cell.text_frame.paragraphs:
                for run in p.runs:
                    run.font.size = Pt(12)
                    run.font.name = FONT
                    run.font.color.rgb = BODY_TEXT

    _add_footer(slide)
    return slide


async def _llm_fill_plan(deck_type: str, title: str, content, plan_headings: list[str]) -> dict:
    content_str = content if isinstance(content, str) else str(content or {})
    headings_str = ", ".join(plan_headings) if plan_headings else "(infer appropriate sections from the content)"

    prompt = f"""Create content for a {DECK_LABELS.get(deck_type, 'presentation')} slide deck titled "{title}".

CONTENT / DATA PROVIDED:
{content_str}

Return JSON only:
{{"title": "...", "subtitle": "...", "slides": [
  {{"heading": "...", "bullets": ["...", "..."]}},
  {{"heading": "...", "type": "table", "headers": ["..."], "rows": [["...", "..."]]}}
]}}

Use these section headings as a guide, in this order: {headings_str}.
For decks with comparable sales / pricing data, use a "table" slide with headers like Address, Beds/Baths, Sqft, Sold Price, Sold Date.
Keep bullets short (under 15 words), 3-6 bullets per slide."""

    raw = await call_claude(system=_PLAN_SYSTEM, prompt=prompt, model=SONNET, max_tokens=2000)
    parsed = parse_json_response(raw)
    if not parsed or not parsed.get("slides"):
        raise ValueError("model did not return a usable slide plan")
    return parsed


async def generate_presentation(deck_type: str, title: str = "", content=None) -> ConnectorResult:
    deck_type = (deck_type or "custom").strip().lower()
    if deck_type not in SLIDE_PLANS:
        deck_type = "custom"

    title = title or DECK_LABELS[deck_type]

    try:
        structured = await _llm_fill_plan(deck_type, title, content, SLIDE_PLANS[deck_type])
    except Exception as e:
        return ConnectorResult(ok=False, error=f"Couldn't generate deck content: {e}")

    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H

    deck_title = structured.get("title") or title
    _add_title_slide(prs, deck_title, structured.get("subtitle", ""))

    for slide_data in structured.get("slides", []):
        if slide_data.get("type") == "table":
            _add_table_slide(prs, slide_data.get("heading", ""), slide_data.get("headers", []), slide_data.get("rows", []))
        else:
            _add_content_slide(prs, slide_data.get("heading", ""), slide_data.get("bullets", []))

    buf = io.BytesIO()
    prs.save(buf)

    safe_title = "".join(c if c.isalnum() or c in " -_" else "" for c in deck_title)[:40].strip() or "presentation"
    doc = save_document(
        buf.getvalue(),
        f"{safe_title}.pptx",
        "application/vnd.openxmlformats-officedocument.presentationml.presentation",
    )

    return ConnectorResult(
        ok=True,
        data={
            "download_url": doc["download_url"],
            "filename": doc["filename"],
            "title": deck_title,
            "slide_count": len(structured.get("slides", [])) + 1,
        },
    )
