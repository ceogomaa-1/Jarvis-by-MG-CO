"""
Dump Learn — ingestion (Jarvis Personal, Study Mode).

This is where the feature's token economics actually happen: every dumped source
(pdf/docx/pptx/url/youtube/image/text) is parsed down to lean, lossless canonical
text EXACTLY ONCE, here, before any reasoning model ever sees it. Every later step
(skeleton condensation, explain at any comprehension level, re-explain after the
user turns the knob) reads that stored text instead of re-parsing the source —
so a 40-page PDF costs tokens once, not on every request.

Large items additionally get:
  - a Haiku "skeleton" pass (condensed outline) so the explain step doesn't have
    to re-read the full raw text for every request:
  - chunked + embedded into dump_learn_chunks (DELIBERATELY its own table, not
    document_chunks/jarvis_skills — this material must never surface in Business
    or Personal memory retrieval) once a whole bin's combined text crosses a size
    threshold, so a textbook-sized dump is retrieved from, not stuffed whole into
    every call.
"""
from __future__ import annotations

import io
import os
import re

import httpx

from backend.tools.url_fetch import fetch_url_content

SUPABASE_URL = os.getenv("SUPABASE_URL") or os.getenv("NEXT_PUBLIC_SUPABASE_URL", "")
SUPABASE_KEY = os.getenv("SUPABASE_SERVICE_ROLE_KEY", "")
ANTHROPIC_API_KEY = os.getenv("ANTHROPIC_API_KEY", "")

HAIKU = "claude-haiku-4-5-20251001"

# ~4 chars/token is a stable enough estimate for the shrink-o-meter and for
# deciding when an item/bin is "large" — it doesn't need to be exact.
CHARS_PER_TOKEN = 4
SKELETON_THRESHOLD_CHARS = 6_000 * CHARS_PER_TOKEN     # ~6k tokens per item
RETRIEVAL_THRESHOLD_CHARS = 30_000 * CHARS_PER_TOKEN   # ~30k tokens combined per bin

CHUNK_SIZE = 800
CHUNK_OVERLAP = 100

ITEMS_TABLE = "dump_learn_items"
CHUNKS_TABLE = "dump_learn_chunks"
UPLOADS_BUCKET = "dump-learn-uploads"


async def _download_storage_object(path: str) -> bytes | None:
    """Backend-side download from the private dump-learn-uploads bucket via the
    service-role key — the file itself was already put there by a direct client
    upload (same convention as personal-chat-attachments)."""
    if not SUPABASE_URL or not SUPABASE_KEY or not path:
        return None
    try:
        async with httpx.AsyncClient() as client:
            resp = await client.get(
                f"{SUPABASE_URL}/storage/v1/object/{UPLOADS_BUCKET}/{path}",
                headers=_read_headers(),
                timeout=30.0,
            )
        if resp.status_code == 200:
            return resp.content
    except Exception as e:
        print(f"DUMP_LEARN: storage download error: {e}")
    return None


def estimate_tokens(text: str | None) -> int:
    """Cheap, consistent token estimate — drives the shrink-o-meter display."""
    return max(0, round(len(text or "") / CHARS_PER_TOKEN))


def _read_headers() -> dict:
    return {"apikey": SUPABASE_KEY, "Authorization": f"Bearer {SUPABASE_KEY}"}


def _write_headers(prefer: str = "return=representation") -> dict:
    return {**_read_headers(), "Content-Type": "application/json", "Prefer": prefer}


def _extract_text(data: dict) -> str:
    """Join every text-type content block. Never assume content[0] is text —
    other block types can lead or the list can be empty (matches the same
    defensive filter study_routes.py applies via the SDK: `if b.type == 'text'`)."""
    blocks = data.get("content") or []
    return "".join(b.get("text", "") for b in blocks if isinstance(b, dict) and b.get("type") == "text").strip()


# ── Per-kind extraction ───────────────────────────────────────────────────────

def _pdf_text(content: bytes) -> str:
    import pypdf
    reader = pypdf.PdfReader(io.BytesIO(content))
    return "\n".join(page.extract_text() or "" for page in reader.pages)


def _docx_text(content: bytes) -> str:
    import docx
    doc = docx.Document(io.BytesIO(content))
    return "\n".join(p.text for p in doc.paragraphs if p.text.strip())


def _pptx_text(content: bytes) -> str:
    import pptx
    prs = pptx.Presentation(io.BytesIO(content))
    slides = []
    for i, slide in enumerate(prs.slides, 1):
        lines = [f"## Slide {i}"]
        for shape in slide.shapes:
            if getattr(shape, "has_text_frame", False) and shape.text_frame.text.strip():
                lines.append(shape.text_frame.text.strip())
            if getattr(shape, "has_table", False):
                for row in shape.table.rows:
                    lines.append(" | ".join(c.text for c in row.cells))
        if slide.has_notes_slide and slide.notes_slide.notes_text_frame.text.strip():
            lines.append(f"(speaker notes: {slide.notes_slide.notes_text_frame.text.strip()})")
        if len(lines) > 1:
            slides.append("\n".join(lines))
    return "\n\n".join(slides)


_YOUTUBE_ID_RE = re.compile(r"(?:youtu\.be/|youtube\.com/(?:watch\?v=|embed/|shorts/))([\w-]{11})")


def _youtube_id(url: str) -> str | None:
    m = _YOUTUBE_ID_RE.search(url or "")
    return m.group(1) if m else None


async def _youtube_transcript(url: str) -> tuple[str, str | None]:
    """Caption-track transcript only (v1 scope) — no audio download/transcription.
    Returns (text, error)."""
    video_id = _youtube_id(url)
    if not video_id:
        return "", "That doesn't look like a YouTube link (couldn't find a video id)."
    try:
        from youtube_transcript_api import YouTubeTranscriptApi
    except ImportError:
        return "", "Video links aren't supported on this deployment yet (missing dependency)."

    def _fetch() -> str:
        api = YouTubeTranscriptApi()
        fetched = api.fetch(video_id)
        return " ".join(snippet.text for snippet in fetched)

    try:
        import asyncio
        text = await asyncio.to_thread(_fetch)
    except Exception as e:
        name = type(e).__name__
        if "TranscriptsDisabled" in name or "NoTranscriptFound" in name:
            return "", "This video has no captions available to read."
        return "", f"Couldn't read this video's transcript ({name})."
    if not text.strip():
        return "", "This video has no captions available to read."
    return text, None


_IMAGE_TO_TEXT_PROMPT = (
    "Transcribe and describe everything in this image as plain text so it can be "
    "studied from. Transcribe ALL visible text exactly (headings, labels, numbers, "
    "tables, diagrams' captions). Then briefly describe any diagrams/charts shown. "
    "Output only the transcription + description — no preamble."
)


async def _image_to_text(content: bytes, media_type: str) -> tuple[str, str | None]:
    import base64
    if not ANTHROPIC_API_KEY:
        return "", "Image reading isn't configured on this deployment."
    b64 = base64.b64encode(content).decode("ascii")
    try:
        async with httpx.AsyncClient() as client:
            resp = await client.post(
                "https://api.anthropic.com/v1/messages",
                headers={
                    "x-api-key": ANTHROPIC_API_KEY,
                    "anthropic-version": "2023-06-01",
                    "content-type": "application/json",
                },
                json={
                    "model": HAIKU,
                    "max_tokens": 2000,
                    "messages": [{
                        "role": "user",
                        "content": [
                            {"type": "image", "source": {"type": "base64", "media_type": media_type, "data": b64}},
                            {"type": "text", "text": _IMAGE_TO_TEXT_PROMPT},
                        ],
                    }],
                },
                timeout=60.0,
            )
        if resp.status_code != 200:
            return "", f"Couldn't read this image (API error {resp.status_code})."
        text = _extract_text(resp.json())
    except Exception as e:
        return "", f"Couldn't read this image ({type(e).__name__})."
    if not text.strip():
        return "", "Couldn't read anything in this image — it may be blank or too low-resolution."
    return text, None


async def extract_source(
    kind: str,
    *,
    content: bytes | None = None,
    url: str | None = None,
    text: str | None = None,
    media_type: str | None = None,
    filename: str | None = None,
) -> dict:
    """Unified entry point: returns {"text": str, "error": str|None, "label": str}."""
    label = filename or url or "Pasted text"
    try:
        if kind == "text":
            body = (text or "").strip()
            if not body:
                return {"text": "", "error": "That's empty — paste some text to learn from.", "label": label}
            return {"text": body, "error": None, "label": label}

        if kind == "url":
            result = await fetch_url_content(url or "")
            if result.get("error"):
                return {"text": "", "error": f"Couldn't read that link ({result['error']}).", "label": label}
            label = result.get("title") or label
            return {"text": result["content"], "error": None, "label": label}

        if kind == "youtube":
            transcript, err = await _youtube_transcript(url or "")
            if err:
                return {"text": "", "error": err, "label": label}
            return {"text": transcript, "error": None, "label": label}

        if kind == "image":
            transcribed, err = await _image_to_text(content or b"", media_type or "image/jpeg")
            if err:
                return {"text": "", "error": err, "label": label}
            return {"text": transcribed, "error": None, "label": label}

        if kind == "pdf":
            body = _pdf_text(content or b"")
            if not body.strip():
                return {"text": "", "error": "Couldn't find selectable text in this PDF — it may be scanned/image-only.", "label": label}
            return {"text": body, "error": None, "label": label}

        if kind == "docx":
            body = _docx_text(content or b"")
            if not body.strip():
                return {"text": "", "error": "This Word file has no readable text in it.", "label": label}
            return {"text": body, "error": None, "label": label}

        if kind == "pptx":
            body = _pptx_text(content or b"")
            if not body.strip():
                return {"text": "", "error": "Couldn't find any text on these slides.", "label": label}
            return {"text": body, "error": None, "label": label}

        return {"text": "", "error": f"Unsupported source type: {kind}", "label": label}
    except Exception as e:
        return {"text": "", "error": f"Couldn't read this ({type(e).__name__}).", "label": label}


# ── Skeleton pass (Haiku condensation for large items) ────────────────────────

_SKELETON_PROMPT = """\
Condense the material below into a compact study outline: section headings, one \
crisp line per section summarizing it, and a short list of key terms/names/numbers \
worth preserving. Be thorough about coverage (don't skip sections) but terse in \
wording — this outline stands in for the full text in later steps.

SOURCE: {label}
MATERIAL:
{content}

Output plain markdown only — headings + bullets. No preamble, no closing remarks."""


async def build_skeleton(text: str, label: str) -> str:
    if not ANTHROPIC_API_KEY or not text.strip():
        return ""
    try:
        async with httpx.AsyncClient() as client:
            resp = await client.post(
                "https://api.anthropic.com/v1/messages",
                headers={
                    "x-api-key": ANTHROPIC_API_KEY,
                    "anthropic-version": "2023-06-01",
                    "content-type": "application/json",
                },
                json={
                    "model": HAIKU,
                    "max_tokens": 1200,
                    "messages": [{"role": "user", "content": _SKELETON_PROMPT.format(label=label, content=text[:60000])}],
                },
                timeout=45.0,
            )
        if resp.status_code != 200:
            print(f"DUMP_LEARN: skeleton API error {resp.status_code}: {resp.text[:200]}")
            return ""
        return _extract_text(resp.json())
    except Exception as e:
        print(f"DUMP_LEARN: skeleton error: {e}")
        return ""


# ── Chunk + embed (retrieval store for large bins only) ───────────────────────

def _chunk_text(text: str) -> list[str]:
    chunks: list[str] = []
    start = 0
    n = len(text)
    while start < n:
        end = start + CHUNK_SIZE
        chunk = text[start:end]
        if chunk.strip():
            chunks.append(chunk)
        if end >= n:
            break
        start = end - CHUNK_OVERLAP
    return chunks


async def chunk_and_embed_item(item_id: str, user_id: str, text: str) -> int:
    """Chunk + embed one item's text into dump_learn_chunks. Best-effort: returns 0
    (never raises) when there's no embedding backend configured — the reasoning
    pass still works fine off skeleton/raw text, just without retrieval."""
    if not SUPABASE_URL or not SUPABASE_KEY or not text.strip():
        return 0
    try:
        from backend.routes.documents import _embed
    except Exception:
        return 0

    chunks = _chunk_text(text)
    if not chunks:
        return 0
    embeddings = await _embed(chunks)
    if not embeddings:
        return 0

    rows = [
        {"item_id": item_id, "user_id": user_id, "chunk_index": i, "content": c, "embedding": e}
        for i, (c, e) in enumerate(zip(chunks, embeddings))
    ]
    stored = 0
    try:
        async with httpx.AsyncClient() as client:
            for i in range(0, len(rows), 20):
                batch = rows[i:i + 20]
                resp = await client.post(
                    f"{SUPABASE_URL}/rest/v1/{CHUNKS_TABLE}",
                    headers=_write_headers("return=minimal"),
                    json=batch,
                    timeout=20.0,
                )
                if resp.status_code in (200, 201, 204):
                    stored += len(batch)
    except Exception as e:
        print(f"DUMP_LEARN: chunk/embed store error: {e}")
    return stored


async def maybe_chunk_and_embed_bin(bin_id: str, user_id: str) -> None:
    """After an item finishes parsing, check whether the WHOLE bin has crossed the
    retrieval threshold; if so, embed any ready items that don't have chunks yet.
    Small bins never pay for embeddings — the reasoning pass just reads raw text."""
    if not SUPABASE_URL or not SUPABASE_KEY:
        return
    try:
        async with httpx.AsyncClient() as client:
            resp = await client.get(
                f"{SUPABASE_URL}/rest/v1/{ITEMS_TABLE}",
                headers=_read_headers(),
                params={"select": "id,extracted_text,token_estimate", "bin_id": f"eq.{bin_id}", "status": "eq.ready"},
                timeout=15.0,
            )
        if resp.status_code != 200:
            return
        items = resp.json()
    except Exception as e:
        print(f"DUMP_LEARN: bin fetch error: {e}")
        return

    total_chars = sum(len(it.get("extracted_text") or "") for it in items)
    if total_chars < RETRIEVAL_THRESHOLD_CHARS:
        return

    for it in items:
        text = it.get("extracted_text") or ""
        if not text.strip():
            continue
        try:
            async with httpx.AsyncClient() as client:
                existing = await client.get(
                    f"{SUPABASE_URL}/rest/v1/{CHUNKS_TABLE}",
                    headers=_read_headers(),
                    params={"select": "id", "item_id": f"eq.{it['id']}", "limit": "1"},
                    timeout=10.0,
                )
            if existing.status_code == 200 and existing.json():
                continue
        except Exception:
            pass
        await chunk_and_embed_item(it["id"], user_id, text)


# ── Item lifecycle (writes to dump_learn_items) ───────────────────────────────

async def _patch_item(item_id: str, fields: dict) -> None:
    if not SUPABASE_URL or not SUPABASE_KEY:
        return
    try:
        async with httpx.AsyncClient() as client:
            await client.patch(
                f"{SUPABASE_URL}/rest/v1/{ITEMS_TABLE}",
                headers=_write_headers("return=minimal"),
                params={"id": f"eq.{item_id}"},
                json=fields,
                timeout=15.0,
            )
    except Exception as e:
        print(f"DUMP_LEARN: item patch error: {e}")


async def ingest_item_task(
    item_id: str,
    bin_id: str,
    user_id: str,
    kind: str,
    *,
    content: bytes | None = None,
    storage_path: str | None = None,
    url: str | None = None,
    text: str | None = None,
    media_type: str | None = None,
    filename: str | None = None,
) -> None:
    """The whole per-item pipeline, meant to run as a detached asyncio task:
    parse -> store lossless text -> (skeleton if large) -> re-check bin-level
    retrieval threshold. Never raises — failures land on the item as status='error'."""
    await _patch_item(item_id, {"status": "parsing"})

    if content is None and storage_path:
        content = await _download_storage_object(storage_path)
        if content is None:
            await _patch_item(item_id, {"status": "error", "error": "Couldn't fetch the uploaded file — try re-adding it."})
            return

    original_size_bytes = len(content) if content else 0
    result = await extract_source(kind, content=content, url=url, text=text, media_type=media_type, filename=filename)
    if result["error"]:
        await _patch_item(item_id, {"status": "error", "error": result["error"]})
        return

    extracted = result["text"]
    char_count = len(extracted)
    token_estimate = estimate_tokens(extracted)

    skeleton = ""
    if char_count > SKELETON_THRESHOLD_CHARS:
        skeleton = await build_skeleton(extracted, result["label"])

    await _patch_item(item_id, {
        "status": "ready",
        "extracted_text": extracted,
        "skeleton_md": skeleton or None,
        "original_size_bytes": original_size_bytes,
        "raw_char_count": char_count,
        "token_estimate": token_estimate,
        "source_name": result["label"],
        "error": None,
    })

    await maybe_chunk_and_embed_bin(bin_id, user_id)
