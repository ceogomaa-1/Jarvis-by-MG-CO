"""
Branding guard (P1.4 / P2): user-facing deliverables must carry the CLIENT'S own
business name and must NEVER show "MG&CO" (that's our company, not the client's).
"""
import io

import pytest


def _pdf_text(pdf_bytes: bytes) -> str:
    pypdf = pytest.importorskip("pypdf")
    reader = pypdf.PdfReader(io.BytesIO(pdf_bytes))
    return "\n".join(page.extract_text() or "" for page in reader.pages)


def test_branded_pdf_uses_client_name_not_mgco():
    from backend.lib.business.pdf_export import HAS_REPORTLAB, generate_branded_document_pdf
    if not HAS_REPORTLAB:
        pytest.skip("reportlab not installed")
    pdf = generate_branded_document_pdf(
        "Listing Package",
        "527 Front St",
        [{"heading": "Overview", "body": "- Spacious\n- Move-in ready"}],
        footer_note="Review with your brokerage.",
        brand_name="Skyline Realty",
    )
    text = _pdf_text(pdf).upper()
    assert "MG&CO" not in text
    assert "SKYLINE REALTY" in text


def test_branded_pdf_without_brand_has_no_mgco():
    from backend.lib.business.pdf_export import HAS_REPORTLAB, generate_branded_document_pdf
    if not HAS_REPORTLAB:
        pytest.skip("reportlab not installed")
    pdf = generate_branded_document_pdf("X", "Y", [{"heading": "H", "body": "body"}])
    assert "MG&CO" not in _pdf_text(pdf).upper()


def _pptx_text(prs) -> str:
    parts = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if shape.has_text_frame:
                parts.append(shape.text_frame.text)
    return " ".join(parts)


def test_pptx_title_slide_uses_client_name_not_mgco():
    pytest.importorskip("pptx")
    from pptx import Presentation
    from backend.lib.business import pptx_generator as pg

    prs = Presentation()
    prs.slide_width = pg.SLIDE_W
    prs.slide_height = pg.SLIDE_H
    pg._add_title_slide(prs, "Listing Presentation", "527 Front St", brand_name="Skyline Realty")
    text = _pptx_text(prs).upper()
    assert "MG&CO" not in text
    assert "SKYLINE REALTY" in text


def test_pptx_title_slide_without_brand_has_no_mgco():
    pytest.importorskip("pptx")
    from pptx import Presentation
    from backend.lib.business import pptx_generator as pg

    prs = Presentation()
    prs.slide_width = pg.SLIDE_W
    prs.slide_height = pg.SLIDE_H
    pg._add_title_slide(prs, "Listing", "Sub", brand_name="")
    assert "MG&CO" not in _pptx_text(prs).upper()
